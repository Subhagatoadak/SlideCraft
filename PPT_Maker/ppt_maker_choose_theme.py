import sys
import os
import json
from pydantic import BaseModel
# Add the root directory to sys.path if not already there.
root_path = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", ".."))
if root_path not in sys.path:
    sys.path.insert(0, root_path)
import io
from llm_service.llm_generator import generate_llm_response, generate_llm_json
import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# For charts
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

# Pydantic model for JSON output
class SlideEvent(BaseModel):
    content: list[str]

# Theme defaults for when no template is uploaded.
THEME_DEFAULTS = {
    "Default": {"bg_color": None, "font_color": None},
    "Dark": {"bg_color": RGBColor(50, 50, 50), "font_color": RGBColor(255, 255, 255)},
    "Corporate": {"bg_color": RGBColor(240, 240, 240), "font_color": RGBColor(0, 0, 0)},
    "Creative": {"bg_color": RGBColor(255, 228, 196), "font_color": RGBColor(75, 0, 130)},
}

# Set page config for a modern wide layout
st.set_page_config(page_title="SlideCraft Pro", page_icon="📊", layout="wide")

# Inject custom CSS for modern styling.
st.markdown(
    """
    <style>
    /* Overall container styling */
    .reportview-container .main .block-container{
        padding-top: 2rem;
        padding-right: 2rem;
        padding-left: 2rem;
        padding-bottom: 2rem;
    }
    /* Custom header styling */
    h1 {
        font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif;
        color: #2c3e50;
        font-weight: 600;
        font-size: 2.8rem;
    }
    h2, h3 {
        font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif;
        color: #34495e;
    }
    /* Sidebar styling */
    .sidebar .sidebar-content {
        font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif;
    }
    /* Input styling tweaks */
    .stTextInput>div>div>input {
        font-size: 1rem;
        padding: 0.5rem;
    }
    .stNumberInput>div>div>input {
        font-size: 1rem;
    }
    </style>
    """,
    unsafe_allow_html=True
)

# Sidebar instructions
st.sidebar.title("Instructions")
st.sidebar.info(
    "Welcome to the SlideCraft Pro!\n\n"
    "1. Fill in your presentation details on the main page.\n"
    "2. Optionally upload a PPTX template file to extract its design and layouts.\n"
    "3. If no template is uploaded, choose a theme for the whole presentation.\n"
    "4. Optionally add background images for title, sections, and slides.\n"
    "5. You can manually create sections and slides or auto-generate slides using AI.\n"
    "6. For manual slides, add rich content with custom fonts, images, charts, and improvement tips.\n"
    "7. When using foreground images, you can upload multiple images.\n\n"
    "Click 'Generate PPT' when you're ready to download your presentation."
)

# Define slide layout options (indices may vary based on your template)
layout_options = {
    "Title Slide (0)": 0,
    "Title and Content (1)": 1,
    "Section Header (2)": 2,
    "Two Content (3)": 3,
    "Comparison (4)": 4,
    "Title Only (5)": 5,
    "Blank (6)": 6,
    "Content with Caption (7)": 7,
    "Picture with Caption (8)": 8,
    "Title and Vertical Text (9)": 9,
    "Vertical Title and Text (10)": 10
}

# Define chart type options.
chart_type_options = {
    "Column Clustered": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "Bar Clustered": XL_CHART_TYPE.BAR_CLUSTERED,
    "Line": XL_CHART_TYPE.LINE,
    "Pie": XL_CHART_TYPE.PIE,
    "Scatter": XL_CHART_TYPE.XY_SCATTER
}

def create_presentation(presentation_title, description, author, template_file, theme_choice,
                        title_bg_bytes, common_content_bg_bytes, sections_data):
    # Use the uploaded template if provided; otherwise create a blank presentation.
    if template_file is not None:
        prs = Presentation(template_file)
    else:
        prs = Presentation()
    
    # (Optional) Debug: print available layouts.
    for i, layout in enumerate(prs.slide_layouts):
        print(f"Layout {i}: {layout.name}")
    
    # ----------------------------
    # Create Main Title Slide
    # ----------------------------
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    
    if title_bg_bytes:
        bg_stream = io.BytesIO(title_bg_bytes)
        bg = slide.shapes.add_picture(bg_stream, 0, 0,
                                      width=prs.slide_width,
                                      height=prs.slide_height)
        # Move image behind other shapes.
        bg._element.getparent().remove(bg._element)
        slide.shapes._spTree.insert(2, bg._element)
    elif not common_content_bg_bytes and (not template_file) and theme_choice and theme_choice != "Default":
        # Apply theme background to title slide if no background image is provided.
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = THEME_DEFAULTS[theme_choice]["bg_color"]
    
    slide.shapes.title.text = presentation_title
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = f"{description}\n\nAuthor: {author}"
    else:
        txBox = slide.shapes.add_textbox(Inches(1), Inches(2),
                                         prs.slide_width - Inches(2),
                                         Inches(1))
        txBox.text = f"{description}\n\nAuthor: {author}"
    
    # ----------------------------
    # Process Each Section
    # ----------------------------
    for section in sections_data:
        section_title = section["section_title"]
        section_header_bg = section.get("section_header_bg", None)
        slides = section["slides"]
        
        # Create section header slide.
        try:
            section_layout = prs.slide_layouts[2]
        except IndexError:
            section_layout = prs.slide_layouts[0]
        sec_slide = prs.slides.add_slide(section_layout)
        if sec_slide.shapes.title:
            sec_slide.shapes.title.text = section_title
        else:
            txBox = sec_slide.shapes.add_textbox(Inches(1), Inches(1),
                                                 prs.slide_width - Inches(2),
                                                 Inches(1))
            txBox.text = section_title
        
        # Add background for section header if provided.
        if section_header_bg:
            bg_stream = io.BytesIO(section_header_bg)
            bg = sec_slide.shapes.add_picture(bg_stream, 0, 0,
                                              width=prs.slide_width,
                                              height=prs.slide_height)
            bg._element.getparent().remove(bg._element)
            sec_slide.shapes._spTree.insert(2, bg._element)
        elif not common_content_bg_bytes and (not template_file) and theme_choice and theme_choice != "Default":
            fill = sec_slide.background.fill
            fill.solid()
            fill.fore_color.rgb = THEME_DEFAULTS[theme_choice]["bg_color"]
        
        # Create slides for this section.
        for idx, slide_data in enumerate(slides):
            layout_index = slide_data.get("layout", 6)
            content = slide_data.get("content", "")
            image_data = slide_data.get("image", None)  # can be a single value or list
            image_type = slide_data.get("image_type", None)  # "background" or "foreground"
            chart_type = slide_data.get("chart_type", None)
            font_size = slide_data.get("font_size", 24)
            font_type = slide_data.get("font_type", "Calibri")
            improvement_tips = slide_data.get("improvement_tips", "")
            
            try:
                slide_layout = prs.slide_layouts[layout_index]
            except IndexError:
                slide_layout = prs.slide_layouts[6]
            new_slide = prs.slides.add_slide(slide_layout)
            
            # Set a default title for the slide if available.
            if new_slide.shapes.title:
                new_slide.shapes.title.text = f"{section_title} - Slide {idx+1}"
            
            # Add text content and apply font settings.
            if content:
                if len(new_slide.placeholders) > 1:
                    placeholder = new_slide.placeholders[1]
                    placeholder.text = content
                    text_frame = placeholder.text_frame
                else:
                    txBox = new_slide.shapes.add_textbox(Inches(1), Inches(2),
                                                         prs.slide_width - Inches(2),
                                                         Inches(2))
                    txBox.text = content
                    text_frame = txBox.text_frame
                # Apply font formatting.
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(font_size)
                        run.font.name = font_type
                        if (not template_file) and theme_choice and theme_choice != "Default" and THEME_DEFAULTS[theme_choice]["font_color"]:
                            run.font.color.rgb = THEME_DEFAULTS[theme_choice]["font_color"]
            
            # Add images if provided.
            if image_data:
                if image_type == "foreground" and isinstance(image_data, list):
                    margin = Inches(0.5)
                    img_width = Inches(3)
                    x = prs.slide_width - img_width - margin
                    y = prs.slide_height - img_width - margin
                    for img_bytes in image_data:
                        img_stream = io.BytesIO(img_bytes)
                        new_slide.shapes.add_picture(img_stream, x, y, width=img_width)
                        x -= (img_width + Inches(0.2))
                else:
                    if isinstance(image_data, list):
                        image_data = image_data[0]
                    img_stream = io.BytesIO(image_data)
                    pic = new_slide.shapes.add_picture(img_stream, 0, 0,
                                                       width=prs.slide_width,
                                                       height=prs.slide_height)
                    pic._element.getparent().remove(pic._element)
                    new_slide.shapes._spTree.insert(2, pic._element)
            elif not common_content_bg_bytes and (not template_file) and theme_choice and theme_choice != "Default":
                # If no image is provided and no common background, apply theme background.
                fill = new_slide.background.fill
                fill.solid()
                fill.fore_color.rgb = THEME_DEFAULTS[theme_choice]["bg_color"]
            
            # Add a chart if requested.
            if chart_type:
                chart_const = chart_type_options.get(chart_type, None)
                if chart_const:
                    chart_data = CategoryChartData()
                    chart_data.categories = ['A', 'B', 'C']
                    chart_data.add_series('Series 1', (10, 20, 30))
                    x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)
                    new_slide.shapes.add_chart(chart_const, x, y, cx, cy, chart_data)
            
            # Add slide notes with improvement tips.
            try:
                notes_slide = new_slide.notes_slide
            except AttributeError:
                notes_slide = new_slide.notes_slide
            notes_slide.notes_text_frame.text = improvement_tips
    
    ppt_io = io.BytesIO()
    prs.save(ppt_io)
    ppt_io.seek(0)
    return ppt_io

def main():
    st.title("SlideCraft Pro")
    st.write("Configure your presentation details below.")
    
    # --- Basic Presentation Details ---
    col1, col2 = st.columns(2)
    with col1:
        presentation_title = st.text_input("Presentation Title", "My Presentation")
    with col2:
        author = st.text_input("Author", "John Doe")
    description = st.text_area("Description", "This is a description for the presentation.")
    
    st.markdown("---")
    
    # --- PPT Template Upload ---
    ppt_template = st.file_uploader("Upload PPT Template (optional)", type=["pptx"], key="ppt_template")
    if ppt_template is None:
        theme_choice = st.selectbox("Choose a Theme", ["Default", "Dark", "Corporate", "Creative"])
    else:
        theme_choice = None
    
    st.markdown("---")
    
    # --- Title Slide Background ---
    add_title_bg = st.checkbox("Add a background image for the title slide?")
    title_bg_bytes = None
    if add_title_bg:
        title_bg_file = st.file_uploader("Upload title slide background", type=["png", "jpg", "jpeg"], key="title_bg")
        if title_bg_file is not None:
            st.image(title_bg_file, caption="Title Slide Background", use_column_width=True)
            title_bg_bytes = title_bg_file.getvalue()
    
    # --- Common Background for Content Slides ---
    add_common_bg = st.checkbox("Add a common background image for all content slides?")
    common_content_bg_bytes = None
    if add_common_bg:
        common_bg_file = st.file_uploader("Upload common background for slides", type=["png", "jpg", "jpeg"], key="common_bg")
        if common_bg_file is not None:
            st.image(common_bg_file, caption="Common Content Background", use_column_width=True)
            common_content_bg_bytes = common_bg_file.getvalue()
    
    st.markdown("---")
    
    # --- AI Auto-Generation Option ---
    auto_generate = st.checkbox("Auto-generate slides using AI?")
    if auto_generate:
        st.markdown("### Auto-Generate Slides Settings")
        ai_context = st.text_area("Enter AI context for slide generation", "Provide any background or context for the presentation here.")
        ai_prompt = st.text_area("Enter AI prompt for slide generation", "Describe the type of slides or content you need.")
        num_ai_slides = st.number_input("Number of slides to generate", min_value=1, step=1, value=3)
    
    st.markdown("---")
    
    # --- Section Header Background Images ---
    add_section_bg = st.checkbox("Add background images for section header slides?")
    
    # --- Manual Sections & Slides ---
    use_sections = st.checkbox("Manually create sections and slides?")
    sections_data = []
    if use_sections:
        num_sections = st.number_input("Number of Sections", min_value=1, step=1, value=1)
        for s in range(int(num_sections)):
            with st.expander(f"Section {s+1} Details", expanded=True):
                section_title = st.text_input(f"Section {s+1} Title", f"Section {s+1}", key=f"section_title_{s}")
                section_header_bg = None
                if add_section_bg:
                    sec_bg_file = st.file_uploader(f"Upload background for Section {s+1} header", type=["png", "jpg", "jpeg"], key=f"sec_bg_{s}")
                    if sec_bg_file is not None:
                        st.image(sec_bg_file, caption=f"Section {s+1} Header Background", use_column_width=True)
                        section_header_bg = sec_bg_file.getvalue()
                num_slides = st.number_input(f"Number of slides in Section {s+1}", min_value=0, step=1, value=1, key=f"num_slides_{s}")
                add_content = st.checkbox(f"Add content, images, or charts to slides in Section {s+1}?", key=f"add_content_{s}")
                slides = []
                if num_slides > 0:
                    slide_tabs = st.tabs([f"Slide {i+1}" for i in range(int(num_slides))])
                    for i, tab in enumerate(slide_tabs):
                        with tab:
                            layout_choice = st.selectbox(
                                f"Select layout for Slide {i+1}",
                                list(layout_options.keys()),
                                key=f"layout_{s}_{i}"
                            )
                            content = ""
                            image_bytes = None
                            image_type = None
                            chart_type = None
                            use_ai = False
                            ai_prompt_manual = ""
                            font_size = 24  # default
                            font_type = "Calibri"  # default
                            # Improvement tips will be auto-generated by AI below.
                            improvement_tips = ""
                            if add_content:
                                content = st.text_area(f"Content for Slide {i+1}", key=f"content_{s}_{i}")
                                # Checkbox for using AI to rewrite content.
                                use_ai = st.checkbox(f"Use AI to rewrite content for Slide {i+1}?", key=f"use_ai_{s}_{i}")
                                if use_ai:
                                    ai_prompt_manual = st.text_area("Enter AI prompt for rewriting:", key=f"ai_prompt_{s}_{i}")
                                # Options for font size and type.
                                font_size = st.number_input("Font Size", min_value=8, max_value=72, value=24, key=f"font_size_{s}_{i}")
                                font_type = st.selectbox("Font Type", options=["Calibri", "Arial", "Times New Roman", "Verdana", "Comic Sans MS"], key=f"font_type_{s}_{i}")
                                add_slide_image = st.checkbox(f"Add an image for Slide {i+1}?", key=f"add_image_{s}_{i}")
                                if add_slide_image:
                                    image_type = st.radio(f"Image type for Slide {i+1}", options=["background", "foreground"], key=f"img_type_{s}_{i}")
                                    if image_type == "foreground":
                                        slide_image_files = st.file_uploader(f"Upload foreground images for Slide {i+1}", type=["png", "jpg", "jpeg"], key=f"slide_image_{s}_{i}", accept_multiple_files=True)
                                        if slide_image_files:
                                            st.image([f for f in slide_image_files], caption=f"Slide {i+1} Images", use_column_width=True)
                                            image_bytes = [f.getvalue() for f in slide_image_files]
                                    else:
                                        slide_image_file = st.file_uploader(f"Upload background image for Slide {i+1}", type=["png", "jpg", "jpeg"], key=f"slide_image_{s}_{i}")
                                        if slide_image_file is not None:
                                            st.image(slide_image_file, caption=f"Slide {i+1} Image", use_column_width=True)
                                            image_bytes = slide_image_file.getvalue()
                                add_chart = st.checkbox(f"Add a chart to Slide {i+1}?", key=f"add_chart_{s}_{i}")
                                if add_chart:
                                    chart_type = st.selectbox(f"Select chart type for Slide {i+1}",
                                                              list(chart_type_options.keys()),
                                                              key=f"chart_{s}_{i}")
                            slides.append({
                                "layout": layout_options[layout_choice],
                                "content": content,
                                "image": image_bytes,
                                "image_type": image_type,
                                "chart_type": chart_type,
                                "use_ai": use_ai,
                                "ai_prompt": ai_prompt_manual,
                                "font_size": font_size,
                                "font_type": font_type,
                                "improvement_tips": improvement_tips
                            })
                sections_data.append({
                    "section_title": section_title,
                    "section_header_bg": section_header_bg,
                    "slides": slides
                })
    else:
        st.info("No sections selected. A default section with one slide will be added.")
        sections_data.append({
            "section_title": "Default Section",
            "section_header_bg": None,
            "slides": [{
                "layout": layout_options["Title and Content (1)"],
                "content": "",
                "image": None,
                "image_type": None,
                "chart_type": None,
                "use_ai": False,
                "ai_prompt": "",
                "font_size": 24,
                "font_type": "Calibri",
                "improvement_tips": ""
            }]
        })
    
    st.markdown("---")
    if st.button("Generate PPT"):
        # For manual slides, if AI rewriting is requested, update the content.
        for section in sections_data:
            for slide_data in section["slides"]:
                if slide_data.get("use_ai", False):
                    original_content = slide_data.get("content", "")
                    ai_prompt_manual = slide_data.get("ai_prompt", "")
                    if original_content and ai_prompt_manual:
                        slide_data["content"] = generate_llm_response(
                            "Context:\n" + original_content + "\n\n" + "Instructions:\n" + ai_prompt_manual,
                            provider="openai",
                            model="gpt-4o",
                            temperature=0.7
                        )
        # If auto-generation is enabled, override manual sections.
        if auto_generate:
            combined_prompt = (
                f"Context:\n{ai_context}\n\n"
                f"Instructions:\n{ai_prompt}\n\n"
                f"Please generate exactly {num_ai_slides} slide contents for a PowerPoint presentation "
                "as a JSON array of strings. Each string should correspond to the content for one slide. "
                "Do not include any additional text."
            )
            ai_output = generate_llm_json(
                combined_prompt, SlideEvent, provider="openai", model="gpt-4o", temperature=0.7
            )
            print(ai_output)
            try:
                slide_contents = ai_output.content
                if not isinstance(slide_contents, list) or len(slide_contents) != num_ai_slides:
                    raise ValueError("The JSON array does not have the required number of slides.")
            except Exception as e:
                st.error("Error parsing AI output as JSON: " + str(e))
                slide_contents = ["" for _ in range(num_ai_slides)]
            sections_data = [{
                "section_title": "Auto-Generated Slides",
                "section_header_bg": None,
                "slides": [{
                    "layout": layout_options["Title and Content (1)"],
                    "content": slide_contents[i],
                    "image": None,
                    "image_type": None,
                    "chart_type": None,
                    "use_ai": False,
                    "ai_prompt": "",
                    "font_size": 24,
                    "font_type": "Calibri",
                    "improvement_tips": ""  # will be auto-generated below
                } for i in range(num_ai_slides)]
            }]
        # --- Auto-generate Improvement Tips for every slide ---
        for section in sections_data:
            for slide_data in section["slides"]:
                slide_content = slide_data.get("content", "").strip()
                if slide_content:
                    improvement = generate_llm_response(
                        "Based on the following slide content, provide improvement tips to enhance clarity, engagement, and design:\n" + slide_content,
                        provider="openai",
                        model="gpt-4o",
                        temperature=0.7
                    )
                    slide_data["improvement_tips"] = improvement
                else:
                    slide_data["improvement_tips"] = "No content provided for improvement tips."
    
        ppt_file = create_presentation(presentation_title, description, author,
                                       ppt_template, theme_choice, title_bg_bytes, common_content_bg_bytes,
                                       sections_data)
        st.success("Presentation generated successfully!")
        st.download_button(
            label="Download PPT",
            data=ppt_file,
            file_name="advanced_generated_presentation.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )

if __name__ == "__main__":
    main()
